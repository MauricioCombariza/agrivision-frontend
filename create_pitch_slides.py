#!/usr/bin/env python3
"""
AgriVision — El Médico · Pitch CEIS 2026
Genera pitch-medico-slides.pptx con identidad visual oficial de AgriVision
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# ── Rutas de logos ────────────────────────────────────────────────────────────
LOGO_FULL  = "/mnt/c/Mauricio/Agrivision/logos/logo_sin_fondo_720px.png"   # 720×134 RGBA
LOGO_ICON  = "/mnt/c/Mauricio/Agrivision/logos/icono_transparente_256px.png"  # 256×256 RGBA

# ── Dimensiones del slide (16:9 widescreen) ───────────────────────────────────
W = Inches(13.33)
H = Inches(7.5)

# ── Paleta oficial AgriVision 2026 ────────────────────────────────────────────
class C:
    # Fondos
    bg_page    = RGBColor(0x0D, 0x1B, 0x2A)  # Fondo Página — background principal
    bg_section = RGBColor(0x0A, 0x15, 0x20)  # Fondo Sección — headers/footers
    bg_card    = RGBColor(0x16, 0x20, 0x32)  # Fondo Card — contenedores
    bg_card2   = RGBColor(0x1C, 0x2B, 0x3D)  # Fondo Card Alt — filas alternas

    # Verdes primarios
    verde      = RGBColor(0x1B, 0x6B, 0x3A)  # Verde Primario — headers
    verde_med  = RGBColor(0x3D, 0xBA, 0x6E)  # Verde Medio — acentos, íconos
    verde_osc  = RGBColor(0x1E, 0x4D, 0x35)  # Verde Oscuro — cards hover
    verde_txt  = RGBColor(0x7F, 0xDB, 0xA0)  # Verde Texto — sobre fondos oscuros
    verde_ok   = RGBColor(0x2E, 0xCC, 0x71)  # Verde Éxito — confirmaciones OK

    # Dorados / acento
    dorado     = RGBColor(0xF3, 0x9C, 0x12)  # Dorado — highlights, CTAs
    dorado_s   = RGBColor(0xE8, 0xA0, 0x20)  # Dorado suave
    amarillo   = RGBColor(0xF1, 0xC4, 0x0F)  # Amarillo — badges de fase

    # Textos
    white      = RGBColor(0xFF, 0xFF, 0xFF)  # Blanco — títulos
    gris_c     = RGBColor(0xB0, 0xC4, 0xD8)  # Gris Claro — texto secundario
    gris_m     = RGBColor(0x6A, 0x8A, 0xAD)  # Gris Medio — metadatos

    # Semánticos
    rojo       = RGBColor(0xE7, 0x4C, 0x3C)  # Rojo Error — diagnóstico crítico
    naranja    = RGBColor(0xE6, 0x7E, 0x22)  # Naranja Alerta
    azul       = RGBColor(0x2E, 0x86, 0xC1)  # Azul Info
    purpura    = RGBColor(0x7D, 0x3C, 0x98)  # Púrpura — prescripción/números

# ── Helpers de geometría ──────────────────────────────────────────────────────

def p(fx, fy, fw, fh):
    """Fracciones del slide → (x, y, w, h) en Emu."""
    return Emu(int(W * fx)), Emu(int(H * fy)), Emu(int(W * fw)), Emu(int(H * fh))

def bg(slide, color):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color

def rect(slide, fx, fy, fw, fh, color, alpha_rgb=None):
    x, y, w, h = p(fx, fy, fw, fh)
    shape = slide.shapes.add_shape(1, x, y, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = alpha_rgb if alpha_rgb else color
    shape.line.fill.background()
    return shape

def txt(slide, text, fx, fy, fw, fh,
        size=28, color=None, bold=False, italic=False,
        align=PP_ALIGN.CENTER, valign="middle",
        font="Poppins"):
    if color is None:
        color = C.white
    x, y, w, h = p(fx, fy, fw, fh)
    txb = slide.shapes.add_textbox(x, y, w, h)
    tf = txb.text_frame
    tf.word_wrap = True

    from pptx.enum.text import MSO_ANCHOR
    tf.vertical_anchor = {"top": MSO_ANCHOR.TOP,
                          "middle": MSO_ANCHOR.MIDDLE,
                          "bottom": MSO_ANCHOR.BOTTOM}.get(valign, MSO_ANCHOR.MIDDLE)

    lines = text.split("\n")
    p0 = tf.paragraphs[0]
    p0.alignment = align

    for i, line in enumerate(lines):
        para = p0 if i == 0 else tf.add_paragraph()
        para.alignment = align
        run = para.add_run()
        run.text = line
        run.font.size = Pt(size)
        run.font.color.rgb = color
        run.font.bold = bold
        run.font.italic = italic
        run.font.name = font

def logo_icon(slide, fx, fy, size_in=0.55):
    """Ícono cuadrado de AgriVision."""
    x, y = Emu(int(W * fx)), Emu(int(H * fy))
    s = Inches(size_in)
    slide.shapes.add_picture(LOGO_ICON, x, y, s, s)

def logo_full(slide, fx, fy, width_in=4.5):
    """Logo completo (wordmark) con fondo transparente."""
    x, y = Emu(int(W * fx)), Emu(int(H * fy))
    w = Inches(width_in)
    h = Inches(width_in / (720 / 134))  # ratio original 720:134
    slide.shapes.add_picture(LOGO_FULL, x, y, w, h)

def timer(slide, time_str):
    """Badge verde con tiempo — esquina superior derecha."""
    rect(slide, 0.815, 0.03, 0.165, 0.115, C.verde_ok)
    txt(slide, f"⏱  {time_str}", 0.815, 0.03, 0.165, 0.115,
        size=17, color=C.bg_page, bold=True, font="Roboto Mono")

def header(slide, title_text, color=None):
    """Barra de header estándar con ícono del logo."""
    if color is None:
        color = C.verde
    rect(slide, 0, 0, 1.0, 0.125, color)
    # Ícono pequeño a la izquierda
    logo_icon(slide, 0.008, 0.012, size_in=0.65)
    # Texto del header
    txt(slide, title_text, 0.085, 0, 0.71, 0.125,
        size=19, bold=True, align=PP_ALIGN.LEFT, valign="middle")

def footer_brand(slide):
    """Franja inferior mínima con tagline de marca."""
    rect(slide, 0, 0.955, 1.0, 0.045, C.bg_section)
    txt(slide, "DATOS DEL CAMPO. DECISIONES ACERTADAS.  ·  combariza.com",
        0.05, 0.955, 0.90, 0.045,
        size=8, color=C.gris_m, font="Roboto Mono")

def notes(slide, text):
    slide.notes_slide.notes_text_frame.text = text

# ── Diapositivas ──────────────────────────────────────────────────────────────

def slide_portada(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg(s, C.bg_page)

    # Franja superior verde
    rect(s, 0, 0, 1.0, 0.05, C.verde)
    # Línea dorada inferior
    rect(s, 0, 0.95, 1.0, 0.05, C.verde)
    rect(s, 0, 0.93, 1.0, 0.008, C.dorado)

    # Logo grande centrado (top area)
    logo_full(s, fx=0.5 - (4.5 / 13.33) / 2, fy=0.09, width_in=4.5)

    # Título principal
    txt(s, "EL MÉDICO", 0.05, 0.30, 0.90, 0.38,
        size=96, color=C.rojo, bold=True)

    # Subtítulo
    txt(s, "Pitch CEIS 2026  ·  5 minutos",
        0.05, 0.70, 0.90, 0.12,
        size=26, color=C.gris_c)

    # Tagline dorada
    txt(s, '«Porque "más o menos" no es un diagnóstico.»',
        0.05, 0.82, 0.90, 0.10,
        size=18, color=C.dorado, italic=True)

    notes(s, "PORTADA — preparar antes de empezar.\nRespira. Pausa. Luego abre con voz de médico.")
    return s


def slide_paciente(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg(s, C.bg_page)
    header(s, "🏥   HISTORIA CLÍNICA — El paciente", color=C.bg_section)

    # Panel izquierdo — datos del paciente
    rect(s, 0.03, 0.155, 0.44, 0.775, C.bg_card)
    txt(s, "Nombre:\nFloricultora colombiana\n\n200 ha · Sabana Norte\n\nExporta:\nMiami · Holanda · Japón\n\nFactura $100M+ COP/mes",
        0.05, 0.17, 0.40, 0.75,
        size=19, color=C.white, align=PP_ALIGN.LEFT, valign="top")

    # Panel derecho — síntoma
    rect(s, 0.54, 0.155, 0.435, 0.775, C.bg_card2)
    txt(s, "Síntoma principal:", 0.56, 0.175, 0.40, 0.09,
        size=14, color=C.gris_m, align=PP_ALIGN.LEFT, valign="middle")
    txt(s, '«Más o menos\n45.000 flores,\njefe.»',
        0.55, 0.265, 0.42, 0.35,
        size=30, color=C.dorado, bold=True, italic=True)
    txt(s, "Con ese dato confirma pedidos\nde $100M COP a Miami.\nEn 2026. Con IA disponible.",
        0.55, 0.635, 0.42, 0.25,
        size=15, color=C.gris_c, align=PP_ALIGN.CENTER)

    timer(s, "0:45")
    footer_brand(s)
    notes(s, (
        "Vengo a presentarles un paciente.\n"
        "Nombre: floricultora colombiana promedio. 200 hectáreas.\n"
        "Exporta a Miami, Holanda, Japón. Factura cientos de millones.\n"
        "Síntoma principal: toma decisiones de cien millones basándose en la sensación del agrónomo.\n"
        "Y lo más preocupante: el paciente no sabe que está enfermo.\n"
        "Cree que así funciona el negocio. Cree que perder el 25% de una temporada es mala suerte."
    ))
    return s


def slide_diagnostico(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg(s, C.bg_page)

    # Bloque rojo dramático (ocupa casi todo el slide)
    rect(s, 0.03, 0.05, 0.94, 0.82, C.rojo)

    # Línea dorada decorativa top
    rect(s, 0.08, 0.06, 0.84, 0.012, C.dorado)

    txt(s, "CEGUERA DE\nDATOS AGUDA",
        0.05, 0.09, 0.90, 0.54,
        size=84, color=C.white, bold=True)

    txt(s, "Fase crónica.",
        0.05, 0.62, 0.90, 0.13,
        size=36, color=C.dorado, italic=True)

    txt(s, "El paciente no sabe que está enfermo.",
        0.05, 0.75, 0.90, 0.10,
        size=20, color=RGBColor(0xFF, 0xBB, 0xBB))

    rect(s, 0.08, 0.86, 0.84, 0.012, C.dorado)

    # Logo icon bottom-right
    logo_icon(s, fx=0.86, fy=0.89, size_in=0.55)

    timer(s, "0:45")
    notes(s, (
        "El diagnóstico es severo:\n"
        "CEGUERA DE DATOS AGUDA. Fase crónica.\n"
        "(pausa larga — decirlo despacio, con peso clínico)\n"
        "El paciente no sabe que está enfermo."
    ))
    return s


def slide_anamnesis(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg(s, C.bg_page)
    header(s, "ANAMNESIS — Historia del sector")

    # Número grande
    txt(s, "USD $2.500M", 0.04, 0.13, 0.90, 0.40,
        size=92, color=C.dorado, bold=True)
    txt(s, "exportados en flores  ·  Colombia 2025  ·  Récord histórico",
        0.04, 0.52, 0.90, 0.11,
        size=21, color=C.white)

    # Tres métricas
    for i, (valor, label, col) in enumerate([
        ("#2",       "Exportador\nmundial", C.verde_ok),
        ("10.000 ha", "cultivadas\nSabana Norte", C.white),
        ("200.000",   "familias\ndependen", C.white),
    ]):
        x = 0.03 + i * 0.325
        rect(s, x, 0.66, 0.30, 0.28, C.bg_card)
        txt(s, valor, x, 0.66, 0.30, 0.16,
            size=28, color=col, bold=True, valign="middle")
        txt(s, label, x, 0.82, 0.30, 0.12,
            size=13, color=C.gris_c, valign="top")

    timer(s, "1:10")
    footer_brand(s)
    notes(s, (
        "Hagamos la anamnesis. ¿Cómo llegó aquí el paciente?\n"
        "Colombia exportó USD $2.500 millones en flores en 2025. Récord histórico.\n"
        "Segundo exportador mundial. 200.000 familias. 10.000 hectáreas cultivadas."
    ))
    return s


def slide_timeline(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg(s, C.bg_page)
    header(s, "⏱   LÍNEA DE TIEMPO — San Valentín · El momento de la verdad",
           color=C.bg_card)

    # Línea del timeline
    rect(s, 0.07, 0.50, 0.86, 0.022, C.gris_m)

    nodos = [
        (0.09, "ENE 20",  "Confirma\npedido",        C.rojo),
        (0.33, "FEB 1",   "Corte y\nexportación",     C.naranja),
        (0.60, "FEB 8",   "Cadena\nde frío",          C.gris_m),
        (0.84, "FEB 14",  "En el\nflorero",           C.verde_ok),
    ]
    for nx, fecha, evento, col in nodos:
        rect(s, nx - 0.025, 0.42, 0.05, 0.17, col)
        txt(s, fecha,  nx - 0.07, 0.24, 0.14, 0.12,
            size=15, color=col, bold=True, align=PP_ALIGN.CENTER, font="Roboto Mono")
        txt(s, evento, nx - 0.07, 0.37, 0.14, 0.11,
            size=12, color=C.gris_c, align=PP_ALIGN.CENTER)
        txt(s, "↑", nx - 0.015, 0.60, 0.03, 0.07,
            size=14, color=col, align=PP_ALIGN.CENTER)

    # Recuadro con el insight crítico
    rect(s, 0.04, 0.72, 0.92, 0.215, C.verde_osc)
    txt(s, "El gerente confirma el 20 de enero cuántas flores tendrá el 1 de febrero\n— cuando todavía están creciendo en el campo.",
        0.06, 0.72, 0.88, 0.215,
        size=19, color=C.dorado, bold=True, valign="middle")

    timer(s, "1:25")
    footer_brand(s)
    notes(s, (
        "Cada año, el momento más crítico es la primera semana de febrero.\n"
        "Las flores colombianas para el 14 de febrero SALEN DE EL DORADO DESDE EL 1 DE FEBRERO.\n"
        "(pausa — este dato nadie lo sabe — es tu credibilidad de campo en una frase)\n"
        "No el 13. El primero.\n"
        "El gerente confirma con datos de HOY cuántas flores tendrá en dos semanas.\n"
        "Y hoy lo hace... con el ojo del agrónomo. Con un Excel. Con fe."
    ))
    return s


def slide_si_senor(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg(s, C.bg_page)
    header(s, "★   DIAGNÓSTICO — Gerente de floricultora...",
           color=C.dorado_s)
    # Override header text color (dorado bg → dark text)
    txt(s, "★   DIAGNÓSTICO — Gerente de floricultora...",
        0.085, 0, 0.71, 0.125,
        size=19, color=C.bg_page, bold=True, align=PP_ALIGN.LEFT, valign="middle")

    sintomas = [
        "¿No sabía cuántas flores tendría el 1 de febrero?",
        "¿Confirmó el pedido a Miami con un «más o menos»?",
        "¿Le apareció botritis y para el viernes ya comió tres surcos?",
        "¿Tiene una secretaria que es la única que sabe dónde va todo?",
        "¿Solo quiere saber con certeza cuántas flores tiene?",
    ]
    for i, stxt in enumerate(sintomas):
        y = 0.145 + i * 0.158
        rect(s, 0.03, y, 0.765, 0.132, C.bg_card)
        txt(s, stxt, 0.05, y, 0.725, 0.132,
            size=17, color=C.white, align=PP_ALIGN.LEFT, valign="middle")
        rect(s, 0.81, y, 0.16, 0.132, C.verde)
        txt(s, "¿Sí señor?", 0.81, y, 0.16, 0.132,
            size=15, bold=True, valign="middle")

    timer(s, "2:15")
    footer_brand(s)
    notes(s, (
        "SECCIÓN DOLORAN — ritmo de comercial de televisión.\n"
        "Cada '¿Sí señor?' esperar que el público responda o asienta.\n"
        "La risa es parte del ritmo. No apresurarlo.\n\n"
        "¿Llamó al agrónomo, él caminó el lote y le dijo 'más o menos 45.000, jefe'?\n"
        "¿Y usted confirmó el pedido a Miami con ese 'más o menos'? — ¿Sí señor?\n"
        "¿El 5 de febrero descubrió que eran 38.000 y el cliente esperaba 50.000? — ¿Sí señor?\n"
        "¿Le apareció botritis un miércoles y para el viernes ya comió tres surcos? — ¿Sí señor?\n"
        "¿Tiene una secretaria que es la única que sabe dónde están todas las cajas? — ¿Sí señor?\n"
        "¿Solo quiere saber con certeza cuántas flores tiene? — ¿Sí señor?"
    ))
    return s


def slide_tiene_tratamiento(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg(s, C.verde)

    txt(s, "Tranquilo.", 0.05, 0.06, 0.90, 0.25,
        size=72, bold=True)
    txt(s, "El cuadro clínico está claro.", 0.05, 0.31, 0.90, 0.17,
        size=36, color=C.verde_txt)

    rect(s, 0.07, 0.52, 0.86, 0.18, C.verde_osc)
    txt(s, "Ceguera de datos aguda  ·  Fase crónica  ·  Con complicaciones logísticas.",
        0.09, 0.52, 0.82, 0.18,
        size=19, color=C.dorado, bold=True, valign="middle")

    txt(s, "Y lo mejor — tiene tratamiento.", 0.05, 0.74, 0.90, 0.18,
        size=44, bold=True)

    logo_icon(s, fx=0.86, fy=0.02, size_in=0.55)
    timer(s, "2:20")
    notes(s, (
        "Tranquilo. El cuadro clínico está claro.\n"
        "(bajar la voz — tono de médico que cierra el diagnóstico)\n"
        "El diagnóstico es: Ceguera de datos aguda. Fase crónica. Con complicaciones logísticas.\n"
        "Y lo mejor — tiene tratamiento."
    ))
    return s


def slide_agrivision(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg(s, C.bg_page)
    header(s, "AgriVision — El Tratamiento")

    # Panel izquierdo — celular
    rect(s, 0.03, 0.155, 0.22, 0.795, C.verde_osc)
    txt(s, "📱", 0.04, 0.22, 0.20, 0.28, size=64, valign="middle")
    txt(s, "El celular que\nel técnico ya\ntiene en el\nbolsillo",
        0.04, 0.55, 0.20, 0.38,
        size=14, color=C.verde_txt, align=PP_ALIGN.CENTER, valign="top")

    # STAR moment — derecha
    txt(s, "Una foto.", 0.29, 0.155, 0.68, 0.22,
        size=66, bold=True, align=PP_ALIGN.LEFT)
    txt(s, "Datos reales.", 0.29, 0.375, 0.68, 0.22,
        size=66, color=C.dorado, bold=True, align=PP_ALIGN.LEFT)
    txt(s, "En menos de 30 segundos.", 0.29, 0.59, 0.68, 0.18,
        size=38, color=C.verde_ok, bold=True, align=PP_ALIGN.LEFT)
    txt(s, "No es hardware costoso. No hay que instalar sensores.\nSolo el celular que el técnico ya tiene.",
        0.29, 0.79, 0.68, 0.14,
        size=16, color=C.gris_c, align=PP_ALIGN.LEFT, valign="top")

    timer(s, "3:00")
    footer_brand(s)
    notes(s, (
        "Buenas noticias. El paciente tiene tratamiento.\n"
        "Yo llevo más de diez años en esos cultivos. Vi la enfermedad desde adentro.\n\n"
        "AgriVision es la plataforma de inteligencia agrícola\n"
        "que convierte UNA FOTO del cultivo en datos accionables — en menos de 30 segundos.\n\n"
        "(levantar el celular)\n\n"
        "UNA FOTO. DATOS REALES. EN MENOS DE 30 SEGUNDOS.\n\n"
        "No es hardware costoso. No hay que instalar sensores en el campo.\n"
        "Solo el celular que el técnico ya tiene en el bolsillo."
    ))
    return s


def slide_flujo(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg(s, C.bg_page)
    header(s, "Cómo funciona — en 30 segundos")

    pasos = [
        ("📷", "1. Foto", "El técnico fotografía\nel lote desde campo", C.azul),
        ("🧠", "2. Análisis", "Visión por computadora:\nbotón · apertura · listo", C.dorado),
        ("📊", "3. Proyección", "7, 14 y 28 días:\ncuántas flores disponibles", C.verde_ok),
    ]
    for i, (icon, titulo, desc, col) in enumerate(pasos):
        x = 0.04 + i * 0.325
        rect(s, x, 0.155, 0.295, 0.67, C.bg_card)
        txt(s, icon, x + 0.005, 0.17, 0.285, 0.22, size=54, valign="middle")
        txt(s, titulo, x + 0.005, 0.415, 0.285, 0.12,
            size=22, color=col, bold=True, valign="middle")
        txt(s, desc, x + 0.005, 0.545, 0.285, 0.26,
            size=15, color=C.gris_c, align=PP_ALIGN.CENTER, valign="top")
        if i < 2:
            txt(s, "→", x + 0.302, 0.38, 0.025, 0.14,
                size=26, color=C.gris_m, align=PP_ALIGN.CENTER, valign="middle")

    # Resultado
    rect(s, 0.04, 0.84, 0.92, 0.115, C.verde_osc)
    txt(s, "Resultado: el gerente sabe el 20 de enero cuántas rosas puede comprometer para el 1 de febrero.",
        0.06, 0.84, 0.88, 0.115,
        size=17, color=C.dorado, bold=True, valign="middle")

    timer(s, "3:20")
    footer_brand(s)
    notes(s, (
        "El técnico de campo fotografía el lote.\n"
        "Visión por computadora identifica el estado fisiológico de cada planta:\n"
        "botón, apertura, listo para corte.\n"
        "El modelo proyecta cuántas flores en 7, 14 y 28 días.\n"
        "El gerente sabe — a finales de enero — exactamente cuántas rosas puede comprometer."
    ))
    return s


def slide_ceis(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg(s, C.bg_page)
    header(s, "El CEIS — El Especialista que Faltaba", color=C.azul)

    # Panel antes
    rect(s, 0.03, 0.155, 0.44, 0.745, C.bg_card)
    txt(s, "Antes del CEIS", 0.03, 0.155, 0.44, 0.10,
        size=18, color=C.rojo, bold=True, valign="middle")
    txt(s, "✓  Tenía la cura\n✗  No sabía cuánto cobrar\n✗  No sabía a quién primero\n✗  Suposiciones con fe,\n    no con ciencia",
        0.05, 0.27, 0.40, 0.60,
        size=18, color=C.white, align=PP_ALIGN.LEFT, valign="top")

    # Panel después
    rect(s, 0.53, 0.155, 0.44, 0.745, C.verde_osc)
    txt(s, "Con el CEIS", 0.53, 0.155, 0.44, 0.10,
        size=18, color=C.verde_ok, bold=True, valign="middle")
    txt(s, "✓  Business + Lean Canvas\n✓  Cliente ideal identificado\n✓  LTV : CAC  →  500 : 1\n✓  Canal más eficiente:\n    cero pesos",
        0.55, 0.27, 0.40, 0.60,
        size=18, color=C.white, align=PP_ALIGN.LEFT, valign="top")

    txt(s, '«El CEIS no me recetó capital. Me recetó claridad.»',
        0.05, 0.91, 0.90, 0.08,
        size=19, color=C.dorado, bold=True, italic=True)

    timer(s, "4:05")
    footer_brand(s)
    notes(s, (
        "Tener el medicamento no es suficiente para construir un negocio.\n"
        "Cuando llegué al CEIS, tenía la cura.\n"
        "Lo que no tenía claro: cuánto cobrar, a quién primero, cómo demostrar que funcionaba.\n"
        "El CEIS me hizo la historia clínica del negocio.\n"
        "Descubrí que tenía suposiciones que eran más fe que ciencia. (esperar risa)\n"
        "LTV:CAC supera 500 a 1. Canal más eficiente: cero pesos.\n"
        "El CEIS no me recetó capital. Me recetó claridad."
    ))
    return s


def slide_pronostico(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg(s, C.bg_page)
    header(s, "⚠   PRONÓSTICO SIN TRATAMIENTO", color=C.rojo)

    for i, (flag, pais) in enumerate([("🇰🇪", "Kenia"), ("🇳🇱", "Holanda"), ("🇮🇱", "Israel")]):
        x = 0.05 + i * 0.315
        rect(s, x, 0.17, 0.28, 0.38, C.bg_card2)
        txt(s, flag, x, 0.19, 0.28, 0.22, size=50, valign="middle")
        txt(s, pais, x, 0.40, 0.28, 0.12, size=22, color=C.gris_c, bold=True, valign="middle")

    txt(s, "Ya están operando con datos de precisión.",
        0.05, 0.60, 0.90, 0.09,
        size=22, color=C.white)

    rect(s, 0.05, 0.71, 0.90, 0.235, RGBColor(0x2A, 0x08, 0x08))
    txt(s, "Colombia tiene una ventana.\nNo es eterna.",
        0.07, 0.71, 0.86, 0.235,
        size=40, color=C.dorado, bold=True, valign="middle")

    timer(s, "4:20")
    footer_brand(s)
    notes(s, (
        "Si el paciente no se trata, el pronóstico es conocido:\n"
        "Los competidores con datos van a ganar los contratos de exportación.\n"
        "Los compradores en Miami van a exigir confirmaciones precisas, no estimaciones.\n"
        "Ya está pasando en Kenia. En Holanda. En Israel.\n"
        "(pausa)\n"
        "Colombia tiene una ventana. No es eterna."
    ))
    return s


def slide_receta(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg(s, C.bg_page)
    header(s, "💊   PRESCRIPCIÓN — Los números del tratamiento", color=C.purpura)

    # LTV:CAC — la métrica estrella
    rect(s, 0.03, 0.155, 0.44, 0.385, C.bg_card)
    txt(s, "LTV : CAC", 0.03, 0.155, 0.44, 0.10,
        size=18, color=C.purpura, bold=True, valign="middle")
    txt(s, "500 : 1", 0.03, 0.255, 0.44, 0.28,
        size=68, color=C.dorado, bold=True, valign="middle")

    # Escala de ingresos
    rect(s, 0.53, 0.155, 0.44, 0.385, C.bg_card)
    for i, (label, val) in enumerate([
        ("1 cliente · 100 ha", "$14–18M COP/mes"),
        ("10 clientes",         "$150M COP/mes"),
        ("30 clientes",         "$450M COP/mes"),
    ]):
        y = 0.19 + i * 0.114
        txt(s, label, 0.55, y, 0.23, 0.10,
            size=14, color=C.gris_c, align=PP_ALIGN.LEFT, valign="middle")
        txt(s, val, 0.775, y, 0.185, 0.10,
            size=16, color=C.verde_ok, bold=True, align=PP_ALIGN.RIGHT,
            valign="middle", font="Roboto Mono")

    # Mercado total
    rect(s, 0.03, 0.575, 0.90, 0.165, C.bg_card2)
    txt(s, "Mercado AgriTech Latinoamérica  ·  USD $10.400M para 2033  ·  +17% anual",
        0.05, 0.575, 0.86, 0.165,
        size=18, color=C.white, valign="middle")

    # Margen + LTV por cliente
    rect(s, 0.03, 0.755, 0.90, 0.18, C.verde_osc)
    txt(s, "Margen bruto  >70%   ·   LTV 36 meses: $540M COP / cliente   ·   CAC  <$1M COP",
        0.05, 0.755, 0.86, 0.18,
        size=17, color=C.verde_txt, bold=True, valign="middle")

    timer(s, "4:45")
    footer_brand(s)
    notes(s, (
        "La prescripción es simple:\n"
        "Un cliente de 100 hectáreas: $14 a $18 millones COP mensuales.\n"
        "Margen bruto >70%. LTV en 36 meses: $540M COP. CAC: <$1M.\n"
        "Diez clientes: $150M al mes. Treinta: $450M al mes.\n"
        "Mercado agritech en Latinoamérica: USD $10.400M para 2033. +17% anual.\n"
        "Colombia puede liderar ese mercado. Con flores. Con datos. Con la Sabana Norte."
    ))
    return s


def slide_alta_medica(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg(s, C.bg_page)
    header(s, "🏥   ALTA MÉDICA — La Visión", color=C.azul)

    rect(s, 0.03, 0.155, 0.90, 0.655, C.verde_osc)

    txt(s, "2030", 0.05, 0.18, 0.35, 0.30,
        size=88, color=C.verde_med, bold=True, align=PP_ALIGN.LEFT, valign="middle")

    txt(s, "El gerente en Tabio\nabre su dashboard\nel 20 de enero",
        0.43, 0.18, 0.48, 0.22,
        size=21, color=C.white, align=PP_ALIGN.LEFT, valign="top")

    txt(s, "87.000 rosas\nlistas para cortar",
        0.05, 0.48, 0.40, 0.26,
        size=24, color=C.verde_ok, bold=True, align=PP_ALIGN.LEFT, valign="top")

    txt(s, "Confirma a Miami\ncon 94% de precisión.",
        0.55, 0.48, 0.36, 0.26,
        size=24, color=C.dorado, bold=True, align=PP_ALIGN.LEFT, valign="top")

    rect(s, 0.03, 0.82, 0.90, 0.135, C.verde)
    txt(s, "No tiene que rezar.     Solo datos.     Solo certeza.     Eso es el alta médica.",
        0.05, 0.82, 0.86, 0.135,
        size=20, bold=True, valign="middle")

    timer(s, "5:05")
    footer_brand(s)
    notes(s, (
        "El objetivo final no es que el paciente deje de estar enfermo.\n"
        "Es que nunca más vuelva a caer.\n"
        "Que en 2030, el gerente de una floricultora en Tabio\n"
        "abra su dashboard el 20 de enero —\n"
        "vea que tiene 87.000 rosas listas para cortar y exportar el 1 de febrero —\n"
        "confirme el pedido a Miami con 94% de precisión —\n"
        "y no tenga que rezar.\n"
        "Solo datos. Solo certeza. Eso es el alta médica."
    ))
    return s


def slide_cierre(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg(s, C.bg_page)

    rect(s, 0.03, 0.04, 0.90, 0.27, C.verde_osc)
    txt(s, "Si conoces un gerente cuya última temporada salió mal\n→  háganme la referencia.",
        0.06, 0.04, 0.84, 0.27,
        size=22, color=C.white, align=PP_ALIGN.LEFT, valign="middle")

    rect(s, 0.03, 0.335, 0.90, 0.195, C.bg_card)
    txt(s, "Si el CEIS tiene un programa para startups con pacientes reales\ny tratamiento probado  →  quiero ser el primer caso clínico.",
        0.06, 0.335, 0.84, 0.195,
        size=20, color=C.white, align=PP_ALIGN.LEFT, valign="middle")

    # Tagline final — bloque sólido
    rect(s, 0, 0.555, 1.0, 0.445, C.verde)
    txt(s, "AgriVision.", 0.05, 0.57, 0.90, 0.22,
        size=76, bold=True)

    # Logo pequeño junto al nombre
    logo_icon(s, fx=0.82, fy=0.575, size_in=0.80)

    txt(s, '«Porque "más o menos" no es un diagnóstico.»',
        0.05, 0.80, 0.90, 0.17,
        size=28, color=C.dorado, bold=True, italic=True)

    timer(s, "5:20")
    notes(s, (
        "El diagnóstico está hecho. El tratamiento existe.\n\n"
        "Si hay alguien en esta sala que conoce a un gerente de finca\n"
        "que la última temporada le salió mal — háganme la referencia.\n\n"
        "Si el CEIS tiene un programa de aceleración para startups\n"
        "con pacientes reales y tratamiento probado — quiero ser el primer caso clínico.\n\n"
        "(pausa — sonrisa)\n"
        "Segundas opiniones bienvenidas.\n"
        "Pero el paciente no puede esperar otra temporada mala.\n\n"
        "(pausa final — mirar al público)\n"
        "AgriVision. Porque 'más o menos' no es un diagnóstico.\n\n"
        "★★★  SILENCIO. NO LLENAR.  ★★★"
    ))
    return s

# ── Main ──────────────────────────────────────────────────────────────────────

BUILDERS = [
    slide_portada,
    slide_paciente,
    slide_diagnostico,
    slide_anamnesis,
    slide_timeline,
    slide_si_senor,
    slide_tiene_tratamiento,
    slide_agrivision,
    slide_flujo,
    slide_ceis,
    slide_pronostico,
    slide_receta,
    slide_alta_medica,
    slide_cierre,
]

def main():
    prs = Presentation()
    prs.slide_width  = W
    prs.slide_height = H

    print("Generando presentación AgriVision — El Médico...")
    for i, fn in enumerate(BUILDERS, 1):
        fn(prs)
        print(f"  [{i:02d}/{len(BUILDERS)}] {fn.__name__.replace('slide_', '')}")

    out = "pitch-medico-slides.pptx"
    prs.save(out)
    print(f"\n✅  Listo: {out}")
    print(f"   {len(BUILDERS)} diapositivas · 5 minutos · identidad AgriVision 2026")
    print(f"\n   Subir a Google Slides:")
    print(f"   drive.google.com → Nuevo → Subir archivo → {out}")

if __name__ == "__main__":
    main()
